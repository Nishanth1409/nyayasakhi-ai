"""Generate Week 4 PPTX for NyayaSakhi AI — run: python generate_week4_nyayasakhi.py"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent / "Week4_Integration_Testing_and_Project_Evolution.pptx"


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = subtitle


def add_bullets(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    first = True
    if subtitle:
        p = body.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.bold = True
        first = False
    for line in bullets:
        p = body.paragraphs[0] if first else body.add_paragraph()
        first = False
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
        p.space_after = Pt(6)


def add_two_column_bullets(prs, title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tx.text_frame.text = title
    tx.text_frame.paragraphs[0].font.size = Pt(32)
    tx.text_frame.paragraphs[0].font.bold = True

    def box(x, y, w, h, small_title, items):
        t = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = t.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = small_title
        p.font.bold = True
        p.font.size = Pt(22)
        for it in items:
            p = tf.add_paragraph()
            p.text = it
            p.level = 0
            p.font.size = Pt(16)
            p.space_after = Pt(4)

    box(0.5, 1.2, 4.4, 5.5, left_title, left_items)
    box(5.1, 1.2, 4.4, 5.5, right_title, right_items)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Week 4 — Integration, Testing & Debugging",
        "NyayaSakhi AI · Multilingual legal guidance · Project evolution vs earlier weekly deck",
    )

    add_bullets(
        prs,
        "Agenda",
        [
            "Week 4 engineering goals: integrate modules, validate end-to-end, harden errors",
            "What changed since the Week 3 slides: architecture & product alignment",
            "New capabilities: document understanding pipeline + voice controls",
            "Testing strategy, known limits, and deliverables",
        ],
    )

    add_bullets(
        prs,
        "Week 4 — primary objectives",
        [
            "End-to-end: Home → Chat (streaming) → Guidance → Help → Document upload",
            "Single AI entry point: Supabase Edge Function nyaya-chat (SSE)",
            "Client-side extraction: PDF text (pdfjs-dist) + image OCR (tesseract.js)",
            "Verify multilingual UI, strict reply-language prompts, and helpline UX",
        ],
    )

    add_two_column_bullets(
        prs,
        "Evolution: Week 3 deck vs current codebase (honest alignment)",
        "Week 3 presentation described",
        [
            "Title: “Inheritance Dispute Resolver”",
            "Backend: Python FastAPI + REST",
            "Modular services + retrieval baseline",
            "Guided questionnaire as main flow",
            "Localization: EN / Kannada / Hindi",
        ],
        "Current NyayaSakhi AI implementation",
        [
            "Product: NyayaSakhi AI (streaming legal assistant)",
            "Backend: Supabase Edge (Deno) — not FastAPI in this repo",
            "LLM via gateway; RAG/corpus services not in MVP path yet",
            "Chat + guidance wizard + help + new “Upload document” journey",
            "Broader Indian language set + enforced reply language in prompts",
        ],
    )

    add_bullets(
        prs,
        "Why the deck and code diverged — professional explanation",
        [
            "Weekly slides are planning artifacts; implementation followed sponsor/platform choices (Lovable + Supabase).",
            "Serverless Edge Functions replace a long-running FastAPI VM for AI proxying: simpler ops and secure secret storage.",
            "Generative chat with streaming improves perceived speed vs blocking REST JSON for long answers.",
            "Document upload was added after core chat: users can attach PDFs/photos for plain-language explanation.",
            "Takeaway for examiners: requirements (multilingual, legal aid, safety) stayed; the technical stack evolved.",
        ],
    )

    add_bullets(
        prs,
        "Backend / integration architecture (current)",
        [
            "Browser → POST /functions/v1/nyaya-chat with Supabase publishable key",
            "Edge Function builds system prompt (domain, tone, 181 / legal aid) + forwards streaming completion",
            "Secrets: LOVABLE_API_KEY only on server; never embedded in Vite bundle",
            "Same endpoint powers VoiceChat and DocumentUpload “ask AI” after text extraction",
        ],
    )

    add_bullets(
        prs,
        "New feature — Document upload module",
        [
            "Accepts PDF or image; caps ~10 MB and ~8000 chars excerpt for safety/latency",
            "PDF: extract text with pdfjs (up to 15 pages)",
            "Images: OCR with tesseract.js; language hint from UI BCP-47 map",
            "User message forces Gemini to answer only in the selected Indian language",
            "Optional read-aloud after stream completes; voice on/off toggle",
        ],
    )

    add_bullets(
        prs,
        "Language enforcement & accessibility (recent changes)",
        [
            "Document flow: explicit “reply ONLY in {language}” instruction reduces English bleed-through",
            "Voice toggle on upload screen: stop/start TTS without leaving the flow",
            "Full wizard strings translated (earlier milestone) — consistent UX across modules",
            "Header branding: logo + favicon refresh for trust and recognition",
        ],
    )

    add_bullets(
        prs,
        "Integration & system testing (Week 4)",
        [
            "Happy path: stream tokens in chat; upload small PDF; verify explanation language",
            "Failure paths: 429 rate limit, 402 credits, network loss → toast + recoverable UI",
            "Cross-browser: speech APIs differ; document OCR heavier on low-end phones",
            "Vitest present (npm test); extend unit tests for pure helpers next sprint",
        ],
    )

    add_bullets(
        prs,
        "Results & deliverables",
        [
            "Integrated SPA with five top-level views including upload",
            "Documented delta from older Week 3 narrative to deployed architecture",
            "Manual test log + screenshots recommended for report appendix",
            "Week 5: final report, deployment hardening, demo video, consolidated presentation",
        ],
    )

    add_title_slide(prs, "Thank you", "Questions & feedback — NyayaSakhi AI · Week 4")

    prs.save(OUT)
    print(f"Wrote: {OUT}")


if __name__ == "__main__":
    main()
